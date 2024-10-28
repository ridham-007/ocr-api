from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
import pdfplumber
import pytesseract
from PIL import Image
import docx
import openai
import io
import xlrd
import openpyxl
from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
from dotenv import load_dotenv
import os

# Load environment variables from .env file
load_dotenv()

# Set up OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

app = FastAPI()

# Extract text from PDF
def extract_text_from_pdf(file):
    with pdfplumber.open(file) as pdf:
        text = ''.join(page.extract_text() for page in pdf.pages)
    return text

# Extract text from DOCX
def extract_text_from_docx(file):
    doc = docx.Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])

# Extract text from images using OCR
def extract_text_from_image(file):
    image = Image.open(file)
    return pytesseract.image_to_string(image)

# Extract text from XLS/XLSX
def extract_text_from_excel(file):
    try:
        # Open XLS files
        workbook = xlrd.open_workbook(file_contents=file.read())
        sheets = workbook.sheets()
    except:
        # Try to open as XLSX
        file.seek(0)
        workbook = openpyxl.load_workbook(file)
        sheets = workbook.worksheets

    extracted_text = []
    for sheet in sheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) for cell in row if cell is not None])
            extracted_text.append(row_text)
    return "\n".join(extracted_text)

# Extract text from PPTX
def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Extract text from EPUB
def extract_text_from_epub(file):
    book = epub.read_epub(file)
    text = []
    for item in book.get_items():
        if item.get_type() == epub.EpubHtml:
            soup = BeautifulSoup(item.get_body_content(), 'html.parser')
            text.append(soup.get_text())
    return "\n".join(text)

# Extract text from TXT
def extract_text_from_txt(file):
    return file.read().decode('utf-8')

# Summarize text using OpenAI GPT
def summarize_text(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Summarize the following text:\n{text}"}
            ],
            max_tokens=150,
            n=1,
            temperature=0.5,
        )
        return response.choices[0].message['content'].strip()
    except Exception as e:
        return str(e)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    file_type = file.content_type
    text = ""

    try:
        if file_type == 'application/pdf':
            text = extract_text_from_pdf(file.file)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            text = extract_text_from_docx(file.file)
        elif file_type.startswith('image/'):
            text = extract_text_from_image(file.file)
        elif file_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            text = extract_text_from_excel(file.file)
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            text = extract_text_from_pptx(file.file)
        elif file_type == 'application/epub+zip':
            text = extract_text_from_epub(file.file)
        elif file_type == 'text/plain':
            text = extract_text_from_txt(file.file)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        # Summarize the extracted text
        summary = summarize_text(text)

        return JSONResponse(content={'summary': summary, 'full_text': text})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
